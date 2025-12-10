"""
Comprehensive File Parser with Multi-Format Support
Handles: PDF, DOCX, Excel, CSV, Code files, Text files, Presentations, and more.
"""

import os
import fitz  # PyMuPDF - fast PDF parsing
import docx
import csv
import json
import xml.etree.ElementTree as ET


def is_supported_file(file_path: str) -> bool:
    """
    Check if file type is supported for parsing.
    
    Supports:
    - Documents: PDF, DOCX, TXT, MD
    - Spreadsheets: XLSX, XLS, CSV
    - Code files: PY, JS, JAVA, C, CPP, CS, GO, RS, etc.
    - Data formats: JSON, XML, YAML, TOML, INI
    - Web: HTML, CSS, JS
    - Presentations: PPTX
    
    Args:
        file_path: Path to file
        
    Returns:
        True if supported, False otherwise
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()
    
    # Comprehensive list of supported extensions
    supported = {
        # Documents
        '.txt', '.md', '.markdown', '.rst', '.pdf', '.docx',
        
        # Spreadsheets & Data
        '.csv', '.tsv', '.xlsx', '.xls',
        
        # Code files
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.c', '.cpp', '.cc', '.cxx',
        '.h', '.hpp', '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala',
        '.r', '.sql', '.sh', '.bash', '.zsh', '.bat', '.ps1', '.m', '.mm',
        
        # Web
        '.html', '.htm', '.css', '.scss', '.sass', '.less', '.xml', '.svg',
        
        # Config & Data
        '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.config',
        '.env', '.properties',
        
        # Other text formats
        '.log', '.tex', '.bib', '.vim', '.gitignore', '.dockerignore',
        
        # Presentations
        '.pptx'
    }
    
    return extension in supported


def get_file_content(file_path: str) -> str | None:
    """
    Extract text content from file using format-specific parsers.
    
    Args:
        file_path: Path to file
        
    Returns:
        Extracted text content, or None if error or unsupported
    """
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return None
        
    # OOM PROTECTION: Check file size before reading
    try:
        file_size = os.path.getsize(file_path)
        # Limits: 10MB for text/code, 50MB for PDFs/Docs (since they are compressed)
        MAX_TEXT_SIZE = 10 * 1024 * 1024  # 10 MB
        MAX_DOC_SIZE = 50 * 1024 * 1024   # 50 MB
        
        _, extension = os.path.splitext(file_path)
        extension = extension.lower()
        
        # Stricter limit for plain text/code (read into memory directly)
        if extension in {'.txt', '.md', '.py', '.js', '.json', '.xml', '.csv', '.log'}:
            if file_size > MAX_TEXT_SIZE:
                print(f"Skipping {file_path}: File too large ({file_size / 1024 / 1024:.2f} MB > 10 MB limit)")
                return f"[Error: File too large to read directly ({file_size / 1024 / 1024:.2f} MB)]"
                
        # Higher limit for documents (parsed selectively)
        if file_size > MAX_DOC_SIZE:
             print(f"Skipping {file_path}: Document too large ({file_size / 1024 / 1024:.2f} MB > 50 MB limit)")
             return f"[Error: Document too large to parse ({file_size / 1024 / 1024:.2f} MB)]"
             
    except Exception as e:
        print(f"Error checking file size: {e}")
        return None
    
    try:
        # Documents
        if extension == ".pdf":
            return _extract_pdf(file_path)
        elif extension == ".docx":
            return _extract_docx(file_path)
        elif extension == ".pptx":
            return _extract_pptx(file_path)
        
        # Spreadsheets
        elif extension in {'.xlsx', '.xls'}:
            return _extract_excel(file_path)
        elif extension in {'.csv', '.tsv'}:
            return _extract_csv(file_path, delimiter=',' if extension == '.csv' else '\t')
        
        # Structured data
        elif extension == '.json':
            return _extract_json(file_path)
        elif extension == '.xml':
            return _extract_xml(file_path)
        
        # Plain text & code files (all text-based)
        elif is_supported_file(file_path):
            return _extract_text(file_path)
        
        else:
            print(f"Unsupported file type: {extension}")
            return None

    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return None


def _extract_text(file_path: str) -> str:
    """Extract from plain text files (txt, md, code files, etc.)."""
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            print(f"Error with encoding {encoding}: {e}")
            continue
    
    # Fallback: binary read and decode with errors='ignore'
    try:
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    except Exception as e:
        print(f"Fallback read failed: {e}")
        return ""


def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF (10x faster than pypdf)."""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text += page.get_text() + "\n"
        doc.close()
    except Exception as e:
        print(f"PDF extraction error for {file_path}: {e}")
        return ""
    
    return text.strip()


def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX files."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""


def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations."""
    try:
        # Try importing pptx
        try:
            from pptx import Presentation
        except ImportError:
            print("python-pptx not installed, skipping PPTX parsing")
            return ""
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""


def _extract_excel(file_path: str) -> str:
    """Extract text from Excel files (xlsx, xls)."""
    try:
        # Try openpyxl for .xlsx
        if file_path.endswith('.xlsx'):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text_parts = []
                
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    text_parts.append(f"=== Sheet: {sheet_name} ===")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            text_parts.append(row_text)
                
                wb.close()
                return "\n".join(text_parts).strip()
            except ImportError:
                print("openpyxl not installed, trying pandas")
        
        # Fallback to pandas for both .xlsx and .xls
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                text_parts.append(df.to_string(index=False))
            
            return "\n\n".join(text_parts).strip()
        except ImportError:
            print("pandas not installed, cannot parse Excel")
            return ""
        
    except Exception as e:
        print(f"Excel extraction error: {e}")
        return ""


def _extract_csv(file_path: str, delimiter: str = ',') -> str:
    """Extract text from CSV/TSV files."""
    try:
        text_parts = []
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f, delimiter=delimiter)
            for row in reader:
                row_text = "\t".join(row)
                text_parts.append(row_text)
        
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"CSV extraction error: {e}")
        return ""


def _extract_json(file_path: str) -> str:
    """Extract text from JSON files (pretty-printed)."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            # Return formatted JSON for better readability
            return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"JSON extraction error: {e}")
        # Fallback to raw text
        return _extract_text(file_path)


def _extract_xml(file_path: str) -> str:
    """Extract text from XML files."""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        # Extract all text content
        text_parts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                text_parts.append(elem.tail.strip())
        
        return "\n".join(text_parts).strip()
    except Exception as e:
        print(f"XML extraction error: {e}")
        # Fallback to raw text
        return _extract_text(file_path)


# Backward compatibility
def extract_text(file_path: str) -> str | None:
    """Alias for get_file_content for backward compatibility."""
    return get_file_content(file_path)


if __name__ == "__main__":
    # Test the parser
    import sys
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        
        if is_supported_file(test_file):
            print(f"✓ File type supported: {test_file}")
            content = get_file_content(test_file)
            if content:
                print(f"✓ Extracted {len(content)} characters")
                print(f"\nPreview (first 500 chars):\n{content[:500]}")
            else:
                print("✗ Failed to extract content")
        else:
            print(f"✗ File type not supported: {test_file}")
    else:
        print("Usage: python fileParser.py <file_path>")
        print("\nSupported formats:")
        print("  Documents: PDF, DOCX, TXT, MD, PPTX")
        print("  Spreadsheets: XLSX, XLS, CSV, TSV")
        print("  Code: PY, JS, JAVA, C, CPP, GO, RS, and 30+ more")
        print("  Data: JSON, XML, YAML, TOML, INI")
        print("  Web: HTML, CSS, XML, SVG")